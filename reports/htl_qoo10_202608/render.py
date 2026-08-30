import sys, os
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

FONT = "/usr/share/fonts/opentype/ipafont-gothic/ipag.ttf"
if not os.path.exists(FONT):
    import subprocess
    FONT = subprocess.check_output(["fc-match","IPAGothic","-f","%{file}"]).decode().strip()
SCALE = 110/914400.0   # EMU -> px at 110 dpi
_fc = {}
def fnt(sz, bold=False):
    k=(round(sz),bold)
    if k not in _fc: _fc[k]=ImageFont.truetype(FONT, max(6,round(sz)))
    return _fc[k]

def px(v): return int(round(v*SCALE))

def rgb(c, default=(30,30,30)):
    try:
        if c is None: return default
        v = c.rgb
        return (v[0],v[1],v[2])
    except Exception:
        return default

def wrap(draw, text, font, maxw):
    """Wrap allowing breaks between any CJK chars; words for latin."""
    lines=[]
    for para in text.split("\n"):
        if not para: lines.append(""); continue
        cur=""
        for ch in para:
            t=cur+ch
            if draw.textlength(t, font=font) > maxw and cur:
                lines.append(cur); cur=ch
            else:
                cur=t
        lines.append(cur)
    return lines

def draw_shape_bg(d, sh, x,y,w,h):
    fill=None
    try:
        f=sh.fill
        if f.type is not None and f.type==1:  # solid
            fill=rgb(f.fore_color, None)
    except Exception: pass
    line=None; lw=1
    try:
        if sh.line.fill.type==1:
            line=rgb(sh.line.color, None)
            lw=max(1,px(sh.line.width or 0)) if sh.line.width else 1
    except Exception: pass
    alpha=255
    try:
        a=sh.fill._xPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
        if a is not None: alpha=int(int(a.get('val'))/100000*255)
    except Exception: pass
    if fill or line:
        st = sh.shape_type
        try: nm = sh._element.spPr.prstGeom.get('prst')
        except Exception: nm = None
        if nm=='ellipse':
            d.ellipse([x,y,x+w,y+h], fill=fill, outline=line, width=lw)
        elif nm=='roundRect':
            d.rounded_rectangle([x,y,x+w,y+h], radius=max(3,int(min(w,h)*0.10)), fill=fill, outline=line, width=lw)
        elif fill and alpha<255:
            ov=Image.new("RGBA",(max(1,w),max(1,h)),tuple(list(fill)+[alpha]))
            d._image.paste(ov,(x,y),ov)
        else:
            d.rectangle([x,y,x+w,y+h], fill=fill, outline=line, width=lw)

def draw_tf(d, tf, x,y,w,h, issues, sid, label, inset=True):
    try:
        li=px(tf.margin_left or 0); ri=px(tf.margin_right or 0)
        ti=px(tf.margin_top or 0); bi=px(tf.margin_bottom or 0)
    except Exception:
        li=ri=ti=bi=0
    tx=x+li; tw=max(6, w-li-ri); ty=y+ti; th=max(4, h-ti-bi)
    blocks=[]
    total=0
    for p in tf.paragraphs:
        runs=[r for r in p.runs]
        txt="".join(r.text for r in runs)
        sz=None; bold=False; col=(30,30,30)
        for r in runs:
            if r.font.size: sz=r.font.size.pt; break
        if sz is None:
            sz = p.font.size.pt if p.font.size else 12
        if runs:
            bold=bool(runs[0].font.bold)
            col=rgb(runs[0].font.color,(30,30,30))
        if not txt.strip():
            total += px(Emu(int(sz*12700)))*0.5; blocks.append((None,None,None,None,None,sz)); continue
        f=fnt(sz*SCALE*914400/72/ (914400*SCALE/72) if False else sz*110/72.0, bold)
        lines=wrap(d, txt, f, tw)
        lh = (sz*110/72.0)*1.22
        align = str(p.alignment) if p.alignment is not None else "LEFT"
        blocks.append((lines,f,lh,col,align,sz))
        total += lh*len(lines)
    va = str(tf.vertical_anchor) if tf.vertical_anchor is not None else "TOP"
    cy = ty
    if "MIDDLE" in va: cy = ty + max(0,(th-total)/2)
    elif "BOTTOM" in va: cy = ty + max(0,th-total)
    nreal=sum(1 for b in blocks if b[0] is not None)
    if nreal>1: total += (nreal-1)*(blocks[0][5]*110/72.0)*0.30
    if total > th + 2:
        issues.append(f"slide{sid}: OVERFLOW {label!r} text≈{total:.0f}px > box {th}px  [{(blocks[0][0][0] if blocks and blocks[0][0] else '')[:34]}]")
    for b in blocks:
        lines,f,lh,col,align,sz = b
        if lines is None:
            cy += (sz*110/72.0)*0.5; continue
        for ln in lines:
            lw_ = d.textlength(ln, font=f)
            lx = tx
            if align and "CENTER" in align: lx = tx + (tw-lw_)/2
            elif align and "RIGHT" in align: lx = tx + (tw-lw_)
            d.text((lx,cy), ln, font=f, fill=col)
            cy += lh
        cy += (sz*110/72.0)*0.30

def render(path, outdir):
    prs=Presentation(path)
    SW=px(prs.slide_width); SH=px(prs.slide_height)
    all_issues=[]
    os.makedirs(outdir, exist_ok=True)
    for i,slide in enumerate(prs.slides,1):
        img=Image.new("RGB",(SW,SH),"white"); d=ImageDraw.Draw(img)
        try:
            bg=slide.background.fill
            if bg.type==1: img.paste(rgb(bg.fore_color,(255,255,255)),[0,0,SW,SH]); d=ImageDraw.Draw(img)
        except Exception: pass
        # pptxgenjs writes slide bg via <p:bg>; fallback parse
        el=slide._element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}bg')
        if el is not None:
            srgb=el.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
            if srgb is not None:
                c=srgb.get('val'); img.paste(tuple(int(c[j:j+2],16) for j in (0,2,4)),[0,0,SW,SH]); d=ImageDraw.Draw(img)
        issues=[]
        for sh in slide.shapes:
            try: x,y,w,h = px(sh.left), px(sh.top), px(sh.width), px(sh.height)
            except Exception: continue
            if x<0 or y<0 or x+w>SW+1 or y+h>SH+1:
                issues.append(f"slide{i}: OUT-OF-BOUNDS {sh.shape_type} at ({x},{y},{w}x{h}) slide {SW}x{SH}")
            if sh.has_chart:
                d.rectangle([x,y,x+w,y+h], outline=(180,190,205), width=2)
                d.text((x+10,y+10), "[CHART]", font=fnt(14), fill=(140,150,165))
                continue
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    import io
                    pim=Image.open(io.BytesIO(sh.image.blob)).convert("RGBA").resize((max(1,w),max(1,h)))
                    img.paste(pim,(x,y),pim); d=ImageDraw.Draw(img)
                except Exception as e:
                    d.rectangle([x,y,x+w,y+h], outline=(200,120,120), width=2)
                continue
            if sh.shape_type == MSO_SHAPE_TYPE.TABLE or sh.has_table:
                tb=sh.table
                colw=[px(c.width) for c in tb.columns]
                rowh=[px(r.height) for r in tb.rows]
                cy=y
                tot=sum(rowh)
                if y+tot > SH: issues.append(f"slide{i}: TABLE overflows slide bottom (ends {y+tot} > {SH})")
                for ri,row in enumerate(tb.rows):
                    cx=x
                    for ci,cell in enumerate(row.cells):
                        cw=colw[ci]; ch=rowh[ri]
                        fillc=None
                        try:
                            if cell.fill.type==1: fillc=rgb(cell.fill.fore_color,None)
                        except Exception: pass
                        d.rectangle([cx,cy,cx+cw,cy+ch], fill=fillc, outline=(220,227,234))
                        draw_tf(d, cell.text_frame, cx+4, cy, cw-8, ch, issues, i, f"table r{ri}c{ci}")
                        cx+=cw
                    cy+=rowh[ri]
                continue
            draw_shape_bg(d, sh, x,y,w,h)
            if sh.has_text_frame and sh.text_frame.text.strip():
                draw_tf(d, sh.text_frame, x, y, w, h, issues, i, sh.text_frame.text[:28])
        # shape-vs-shape overlap audit (tables/cards/text blocks that collide)
        boxes=[]
        for sh in slide.shapes:
            try: bx,by,bw,bh = px(sh.left),px(sh.top),px(sh.width),px(sh.height)
            except Exception: continue
            kind = "TABLE" if getattr(sh,'has_table',False) else ("CHART" if sh.has_chart else
                   ("TEXT" if sh.has_text_frame and sh.text_frame.text.strip() else "SHAPE"))
            if kind=="TABLE":
                bh = sum(px(r.height) for r in sh.table.rows)
            lbl = (sh.text_frame.text.strip()[:26] if sh.has_text_frame and sh.text_frame.text.strip() else kind)
            boxes.append((kind,bx,by,bw,bh,lbl))
        for a in range(len(boxes)):
            for b in range(a+1,len(boxes)):
                k1,x1,y1,w1,h1,l1 = boxes[a]; k2,x2,y2,w2,h2,l2 = boxes[b]
                if k1=="SHAPE" and k2=="SHAPE": continue
                ox = min(x1+w1,x2+w2)-max(x1,x2); oy = min(y1+h1,y2+h2)-max(y1,y2)
                if ox>4 and oy>4:
                    # a card intentionally sits behind its own text: skip SHAPE-under-TEXT containment
                    if k1=="SHAPE" and x1<=x2+2 and y1<=y2+2 and x1+w1>=x2+w2-2 and y1+h1>=y2+h2-2: continue
                    if k2=="SHAPE" and x2<=x1+2 and y2<=y1+2 and x2+w2>=x1+w1-2 and y2+h2>=y1+h1-2: continue
                    issues.append(f"slide{i}: OVERLAP {k1}<{l1}> x {k2}<{l2}> by {ox}x{oy}px")
        img.save(f"{outdir}/slide-{i:02d}.png")
        all_issues += issues
    print("\n".join(all_issues) if all_issues else "NO GEOMETRY ISSUES DETECTED")
    print(f"\n{len(all_issues)} issue(s); rendered {len(prs.slides.__iter__.__self__._sldIdLst)} slides to {outdir}/")

render(sys.argv[1], sys.argv[2])

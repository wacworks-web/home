import sys
from PIL import ImageFont, ImageDraw, Image
from pptx import Presentation
from pptx.util import Emu
import subprocess
FONT=subprocess.check_output(["fc-match","IPAGothic","-f","%{file}"]).decode().strip()
d=ImageDraw.Draw(Image.new("RGB",(10,10)))
prs=Presentation(sys.argv[1]); target=int(sys.argv[2])
for i,sl in enumerate(prs.slides,1):
    if i!=target: continue
    for sh in sl.shapes:
        if not getattr(sh,'has_table',False): continue
        t=sh.table
        colw=[Emu(c.width).inches for c in t.columns]
        need=[0.0]*len(colw); worst=['']*len(colw)
        for row in t.rows:
            for ci,cell in enumerate(row.cells):
                tf=cell.text_frame
                txt=tf.text.strip()
                if not txt: continue
                sz=12
                for p in tf.paragraphs:
                    for r in p.runs:
                        if r.font.size: sz=r.font.size.pt; break
                    break
                f=ImageFont.truetype(FONT, max(6,round(sz*110/72)))
                wpx=d.textlength(txt,font=f)/110.0
                li=Emu(tf.margin_left or 0).inches; ri=Emu(tf.margin_right or 0).inches
                total=wpx+li+ri
                if total>need[ci]: need[ci]=total; worst[ci]=txt[:34]
        print(f"table at y={Emu(sh.top).inches:.2f} rowH={Emu(t.rows[1].height).inches:.3f} rows={len(t.rows)} totalH={sum(Emu(r.height).inches for r in t.rows):.2f}")
        for ci,(cw,nd,wt) in enumerate(zip(colw,need,worst)):
            flag="  <-- TOO NARROW" if nd>cw+0.005 else ""
            print(f"  col{ci}: width={cw:.2f}  needs={nd:.2f}{flag}   [{wt}]")
        print("  sum width", round(sum(colw),2))

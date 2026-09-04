# Estimate text extents per shape and flag overflow / out-of-bounds boxes.
import sys, unicodedata
from pptx import Presentation
from pptx.util import Emu, Length

SLIDE_W, SLIDE_H = 13.333, 7.5

def char_em(ch):
    if ch in "　": return 1.0
    return 1.0 if unicodedata.east_asian_width(ch) in ("W", "F", "A") else 0.52

def measure(paras, box_w_in):
    """Return (needed_height_in, worst_single_line_in) for the shape's text."""
    total_h, worst = 0.0, 0.0
    for p in paras:
        runs = [r for r in p.runs]
        if not runs:
            continue
        size = max((r.font.size.pt for r in runs if r.font.size), default=14.0)
        # explicit line spacing if set on the paragraph
        ls = p.line_spacing
        if ls is None:
            line_pt = size * 1.2
        elif isinstance(ls, Length):          # absolute spacing (spcPts) -> EMU-backed Length
            line_pt = ls.pt
        else:                                  # a plain float is a multiple of line height
            line_pt = size * float(ls)
        line_h = line_pt / 72.0
        text = "".join(r.text for r in runs)
        for seg in text.split("\n"):
            w = sum(char_em(c) for c in seg) * size / 72.0
            worst = max(worst, w)
            lines = max(1, -(-int(w * 1000) // max(1, int(box_w_in * 1000)))) if box_w_in > 0 else 1
            total_h += lines * line_h
    return total_h, worst

def walk(shapes, out, prefix=""):
    for sh in shapes:
        if sh.shape_type == 6:  # group
            walk(sh.shapes, out, prefix)
            continue
        out.append(sh)

prs = Presentation(sys.argv[1])
issues = 0
for i, slide in enumerate(prs.slides, 1):
    shapes = []
    walk(slide.shapes, shapes)
    for sh in shapes:
        if sh.left is None:
            continue
        x, y = Emu(sh.left).inches, Emu(sh.top).inches
        w, h = Emu(sh.width).inches, Emu(sh.height).inches
        rot = getattr(sh, "rotation", 0) or 0
        # bounds (skip full-bleed backgrounds)
        if not (w > 13 and h > 7):
            if x < -0.02 or y < -0.02 or x + w > SLIDE_W + 0.02 or y + h > SLIDE_H + 0.02:
                if abs(rot) < 1:
                    print(f"S{i:02d} OUT-OF-BOUNDS  x={x:.2f} y={y:.2f} w={w:.2f} h={h:.2f}  '{(sh.text_frame.text[:28] if sh.has_text_frame else '')}'")
                    issues += 1
        if not sh.has_text_frame or not sh.text_frame.text.strip():
            continue
        if abs(rot) > 1:
            continue  # rotated boxes measured in their own frame
        need, worst = measure(sh.text_frame.paragraphs, w)
        if need > h + 0.06:
            print(f"S{i:02d} TEXT-OVERFLOW  need={need:.2f}\" box_h={h:.2f}\" w={w:.2f}\"  '{sh.text_frame.text[:40].replace(chr(10),'/')}'")
            issues += 1
        elif need > h * 0.90 and h > 0.2:
            print(f"S{i:02d} tight          need={need:.2f}\" box_h={h:.2f}\"  '{sh.text_frame.text[:36].replace(chr(10),'/')}'")
print(f"\n{issues} issue(s)")
